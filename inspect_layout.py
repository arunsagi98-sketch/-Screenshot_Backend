from pptx import Presentation

def inspect_layout(filepath):
    prs = Presentation(filepath)
    slide = prs.slides[5] # Slide 6 (mobile)
    for i, shape in enumerate(slide.shapes):
        if hasattr(shape, "text") and shape.text:
            print(f"--- Text Block {i} ---")
            print(shape.text)

if __name__ == "__main__":
    inspect_layout("PDF_format/PDF_format.pptx")
