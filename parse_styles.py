from pptx import Presentation
import sys

def parse_styles(filepath):
    prs = Presentation(filepath)
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        if slide.background.fill.type == 1: # solid
            bg_color = slide.background.fill.fore_color
            if hasattr(bg_color, 'rgb'):
                print(f"Background Color: {bg_color.rgb}")
            elif hasattr(bg_color, 'theme_color'):
                print(f"Background Theme Color: {bg_color.theme_color}")
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(f"Text block: '{shape.text.strip()}'")
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font = run.font
                        color = "Default"
                        if font.color and hasattr(font.color, "rgb") and font.color.rgb:
                            color = str(font.color.rgb)
                        elif font.color and hasattr(font.color, "theme_color"):
                            color = f"Theme {font.color.theme_color}"
                            
                        size = font.size.pt if font.size else "Default"
                        bold = font.bold
                        name = font.name
                        print(f"  -> '{run.text.strip()}': Font={name}, Size={size}, Color={color}, Bold={bold}")
                        
        if i >= 1: # Only need first 2 slides to get the pattern
            break

if __name__ == "__main__":
    import glob
    files = glob.glob("PDF_format/*.pptx")
    if files:
        parse_styles(files[0])
