import os
from pptx import Presentation

def parse_pptx(filepath):
    if not os.path.exists(filepath):
        print(f"File not found: {filepath}")
        return
        
    prs = Presentation(filepath)
    print(f"Total slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(f"Text: {shape.text.strip()}")
            if hasattr(shape, "image"):
                print(f"Picture found")

if __name__ == "__main__":
    import glob
    files = glob.glob("PDF_format/*.pptx")
    if files:
        parse_pptx(files[0])
    else:
        print("No pptx found")
