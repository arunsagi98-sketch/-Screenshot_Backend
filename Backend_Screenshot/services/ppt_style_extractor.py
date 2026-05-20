import os
import re
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation

# OOXML DrawingML namespace (kept for future targeted XPath lookups)
_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def load_scheme_colors_from_pptx(pptx_path):
    """
    Read accent / text colors from ppt/theme/theme1.xml (_clrScheme).
    Returns dict keyed by scheme tag ('accent1','dk1',...) -> 'RRGGBB' uppercase.
    """
    if not pptx_path or not os.path.exists(pptx_path):
        return {}

    scheme = {}
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            # Decks vary: theme/theme1.xml is the usual case
            names = sorted(
                n for n in zf.namelist() if n.startswith("ppt/theme/") and n.endswith(".xml")
            )
            if not names:
                return {}
            xml_bytes = zf.read(names[0])
        tree = ET.fromstring(xml_bytes)
        clr_scheme = None
        for el in tree.iter():
            if el.tag.endswith("clrScheme"):
                clr_scheme = el
                break

        if clr_scheme is None:
            return {}

        def _color_from_clr_element(clr_elem):
            for child in list(clr_elem):
                tag = child.tag.split("}")[-1]
                if tag == "srgbClr":
                    val = child.get("val") or ""
                    val = "".join(filter(str.isalnum, val))[:6]
                    return val.upper() if len(val) == 6 else None
                if tag == "sysClr":
                    lc = child.get("lastClr") or ""
                    lc = "".join(filter(str.isalnum, lc))[:6]
                    return lc.upper() if len(lc) == 6 else None
            return None

        for child in list(clr_scheme):
            tag_local = child.tag.split("}")[-1]
            if tag_local.endswith("Clr"):  # e.g. accent1Clr is wrong; OOXML uses direct children like accent1
                continue
            if not re.match(r"^[a-zA-Z]+\d*$", tag_local):
                continue
            hex_rgb = _color_from_clr_element(child)
            if hex_rgb:
                scheme[tag_local.lower()] = hex_rgb

        print(f"[PPT THEME] Scheme colors extracted: keys={sorted(scheme.keys())}")
        return scheme
    except Exception as e:
        print(f"[PPT THEME] Could not parse theme XML: {e}")
        return {}


def enrich_styles_with_scheme(styles, pptx_path, scheme_override=None):
    """
    When runs use Theme colors PowerPoint resolves via clrScheme — often missing RGB in pptx extraction.
    We patch title/text/accent hex from clrScheme accents / dk hues.
    """
    scheme = scheme_override if scheme_override is not None else load_scheme_colors_from_pptx(pptx_path)
    if not scheme:
        return styles, scheme

    accent1 = scheme.get("accent1")
    dk1 = scheme.get("dk1")
    dk2 = scheme.get("dk2")

    if dk1 and len(dk1) == 6:
        styles["text_color"] = dk1
    title_src = dk2 or dk1 or accent1
    if title_src and len(title_src) == 6:
        styles["title_color"] = title_src

    styles["accent_color"] = (accent1 if accent1 and len(accent1) == 6 else None) or (
        dk2 if dk2 and len(dk2) == 6 else None
    ) or "6366F1"
    styles.setdefault("subtitle_color", (accent1 if accent1 and len(accent1) == 6 else styles["accent_color"]))

    return styles, scheme


def resolve_run_hex(font, scheme):
    """Return RRGGBB for a FontColor if possible (RGB or theme-mapped via scheme tags)."""
    if not scheme:
        return None
    try:
        color = getattr(font, "color", None)
        if color is None:
            return None

        # RGB literal
        rgb = getattr(color, "rgb", None)
        if rgb is not None:
            color_rgb = getattr(rgb, "rgb", rgb)
            try:
                r, g, b = color_rgb[0], color_rgb[1], color_rgb[2]
                return f"{int(r):02X}{int(g):02X}{int(b):02X}"
            except Exception:
                pass

        ENUM_TO_TAG = {
            "DARK_1": "dk1",
            "LIGHT_1": "lt1",
            "DARK_2": "dk2",
            "LIGHT_2": "lt2",
            "ACCENT_1": "accent1",
            "ACCENT_2": "accent2",
            "ACCENT_3": "accent3",
            "ACCENT_4": "accent4",
            "ACCENT_5": "accent5",
            "ACCENT_6": "accent6",
            "TEXT_1": "dk1",
            "TEXT_2": "dk2",
            "BACKGROUND_1": "lt1",
            "BACKGROUND_2": "lt2",
        }

        tm = getattr(color, "theme_color", None)
        name = getattr(tm, "name", None) if tm is not None else None
        if isinstance(tm, int):
            try:
                from pptx.enum.dml import MSO_THEME_COLOR

                name = MSO_THEME_COLOR(tm).name
            except Exception:
                name = None

        scheme_key = ENUM_TO_TAG.get(str(name))
        if scheme_key and scheme.get(scheme_key):
            return scheme[scheme_key]
    except Exception:
        pass
    return None

def extract_ppt_styles(pptx_path):
    """
    Extract text styling information from a PowerPoint file.
    Returns a dictionary with font, size, color, and bold information.
    """
    if not os.path.exists(pptx_path):
        print(f"[PPT STYLE] File not found: {pptx_path}")
        return None
        
    try:
        prs = Presentation(pptx_path)
        scheme = load_scheme_colors_from_pptx(pptx_path)
        
        # Default styles (fallback values)
        styles = {
            "font_name": "Segoe UI",
            "title_size": 30,
            "text_size": 18,
            "title_bold": True,
            "text_bold": True,
            "title_color": "000000",  # Black - default text color
            "text_color": "000000",    # Black - default text color
            "background_color": "FFFFFF",  # White - default background
            "accent_color": (scheme.get("accent1") if scheme else None) or "6366F1",
        }
        
        # Extract from first slide (title slide)
        if len(prs.slides) > 0:
            slide = prs.slides[0]
            
            # Extract background color - try to get actual background
            background_color_found = False
            try:
                if slide.background.fill.type == 1:  # solid
                    bg_color = slide.background.fill.fore_color
                    if hasattr(bg_color, 'rgb') and bg_color.rgb:
                        styles["background_color"] = f"{bg_color.rgb[0]:02X}{bg_color.rgb[1]:02X}{bg_color.rgb[2]:02X}"
                        background_color_found = True
            except Exception:
                pass  # Keep default
            
            # Extract text styles from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            font = run.font
                            
                            # Get font name
                            if font.name:
                                styles["font_name"] = font.name
                            
                            # Get font size
                            if font.size:
                                size_pt = font.size.pt
                                # Determine if this is likely a title (larger size) or regular text
                                if size_pt >= 24:  # Title-sized text
                                    styles["title_size"] = int(size_pt)
                                else:  # Regular text
                                    styles["text_size"] = int(size_pt)
                            
                            # Get bold status
                            if font.bold is not None:
                                if font.bold:  # If explicitly set to bold
                                    # Check if this looks like a label (ends with colon) or value
                                    text = run.text.strip()
                                    if text.endswith(':'):
                                        styles["title_bold"] = True
                                    else:
                                        styles["text_bold"] = True
                            
                            # Get color - try to get explicit color, otherwise use defaults
                            text_color_found = False
                            try:
                                text = run.text.strip()
                                color_hex = resolve_run_hex(font, scheme)
                                if not color_hex:
                                    fc = getattr(font, "color", None)
                                    if fc and getattr(fc, "rgb", None) is not None:
                                        color_rgb = fc.rgb
                                        try:
                                            r, g, b = color_rgb[0], color_rgb[1], color_rgb[2]
                                            color_hex = f"{int(r):02X}{int(g):02X}{int(b):02X}"
                                        except Exception:
                                            color_hex = None
                                if color_hex:
                                    if text.endswith(':') if text else False:
                                        styles["title_color"] = color_hex
                                    else:
                                        styles["text_color"] = color_hex
                            except Exception:
                                pass

        if scheme:
            styles["accent_color"] = (
                scheme.get("accent1") if scheme.get("accent1") else styles.get("accent_color")
            ) or "6366F1"
            if styles.get("title_color") == "000000" and styles.get("text_color") == "000000":
                styles, _ = enrich_styles_with_scheme(styles, pptx_path, scheme)
        else:
            styles.setdefault("accent_color", "6366F1")
                            
        print(f"[PPT STYLE] Extracted styles: {styles}")
        return styles
        
    except Exception as e:
        print(f"[PPT STYLE ERROR] Failed to extract styles: {e}")
        # Return default styles on error
        return {
            "font_name": "Segoe UI",
            "title_size": 30,
            "text_size": 18,
            "title_bold": True,
            "text_bold": True,
            "title_color": "000000",
            "text_color": "000000",
            "background_color": "FFFFFF",
            "accent_color": "6366F1",
        }

def get_ppt_file_path():
    """Returns the found PPT file path or None."""
    possible_paths = [
        "../PPT_Format/CA01437_Banoffee_English_Banner_Screenshots_Apr'26.pptx",
        "PPT_Format/CA01437_Banoffee_English_Banner_Screenshots_Apr'26.pptx",
        "../../PPT_Format/CA01437_Banoffee_English_Banner_Screenshots_Apr'26.pptx",
        "../../../PPT_Format/CA01437_Banoffee_English_Banner_Screenshots_Apr'26.pptx"
    ]
    for path in possible_paths:
        if os.path.exists(path):
            return path
    return None

def get_ppt_styles():
    """
    Get PPT styles from the standard location.
    """
    path = get_ppt_file_path()
    if path:
        return extract_ppt_styles(path)
    
    print("[PPT STYLE] PPT file not found in any expected location")
    return None

def extract_ppt_assets():
    """
    Extracts Slide 1 background (image2.jpg) and Billiontags logo (image3.png) from the PPTX zip.
    Saves them in an 'extracted_ppt_media' folder in the workspace root.
    """
    import zipfile
    pptx_path = get_ppt_file_path()
    if not pptx_path:
        print("[PPT ASSETS] PPTX file not found to extract assets.")
        return None
        
    extract_dir = "extracted_ppt_media"
    os.makedirs(extract_dir, exist_ok=True)
    
    bg_dest = os.path.join(extract_dir, "image2.jpg")
    logo_dest = os.path.join(extract_dir, "image3.png")
    
    # If already extracted, return paths
    if os.path.exists(bg_dest) and os.path.exists(logo_dest):
        return {"background": bg_dest, "logo": logo_dest}
        
    try:
        with zipfile.ZipFile(pptx_path, 'r') as z:
            # 1. Background image — prefer canonical name, otherwise largest JPEG from ppt/media/.
            bg_target = "ppt/media/image2.jpg"
            if bg_target in z.namelist():
                with open(bg_dest, 'wb') as f:
                    f.write(z.read(bg_target))
                print(f"[PPT ASSETS] Extracted cover background to {bg_dest}")
            elif not os.path.exists(bg_dest):
                media = sorted(
                    (n for n in z.namelist() if n.startswith("ppt/media/") and n.lower().endswith((".jpg", ".jpeg"))),
                    key=lambda n: z.getinfo(n).file_size,
                    reverse=True,
                )
                if media:
                    picked = media[0]
                    with open(bg_dest, 'wb') as f:
                        f.write(z.read(picked))
                    print(f"[PPT ASSETS] Extracted fallback cover ({picked}) -> {bg_dest}")
                
            # 2. Logo — prefer canonical name; otherwise PNG with reasonable footprint
            logo_target = "ppt/media/image3.png"
            if logo_target in z.namelist():
                with open(logo_dest, 'wb') as f:
                    f.write(z.read(logo_target))
                print(f"[PPT ASSETS] Extracted logo to {logo_dest}")
            elif not os.path.exists(logo_dest):
                pngs = sorted(
                    (n for n in z.namelist() if n.startswith("ppt/media/") and n.lower().endswith(".png")),
                    key=lambda n: z.getinfo(n).file_size,
                )
                for candidate in pngs:
                    sz = z.getinfo(candidate).file_size
                    if sz < 3_500_000:
                        with open(logo_dest, 'wb') as f:
                            f.write(z.read(candidate))
                        print(f"[PPT ASSETS] Extracted fallback logo ({candidate}) -> {logo_dest}")
                        break
                
        return {"background": bg_dest, "logo": logo_dest}
    except Exception as e:
        print(f"[PPT ASSETS ERROR] Failed to extract media assets: {e}")
        return None

def extract_campaign_details():
    """
    Extracts text details (Title, Start Date, End Date, Format) from Slide 1 shapes.
    """
    pptx_path = get_ppt_file_path()
    details = {
        "title": "Chop Steakhouse FY26 Banner Campaign",
        "start_date": "09th Apr 2026",
        "end_date": "22nd Apr 2026",
        "format": "Banner"
    }
    
    if not pptx_path:
        return details
        
    try:
        prs = Presentation(pptx_path)
        if len(prs.slides) > 0:
            slide = prs.slides[0]
            # Look at all shapes on slide 1
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    lines = [line.strip() for line in text.split('\n') if line.strip()]
                    for line in lines:
                        if line.startswith("Start Date:"):
                            details["start_date"] = line.replace("Start Date:", "").strip()
                        elif line.startswith("End Date:"):
                            details["end_date"] = line.replace("End Date:", "").strip()
                        elif line.startswith("Format:"):
                            details["format"] = line.replace("Format:", "").strip()
                        elif "Campaign" in line:
                            details["title"] = line.strip()
        print(f"[PPT STYLE] Extracted campaign details: {details}")
        return details
    except Exception as e:
        print(f"[PPT STYLE ERROR] Failed to extract campaign details: {e}")
        return details